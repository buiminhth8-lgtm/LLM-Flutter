"""Document loader - parse Word, PDF, TXT, Markdown, CSV, Excel, HTML files."""

from __future__ import annotations

import csv
import hashlib
import json
import re
from dataclasses import dataclass
from pathlib import Path


@dataclass
class Document:
    """A parsed document chunk."""

    content: str
    metadata: dict

    def __repr__(self):
        src = self.metadata.get("source", "unknown")
        return f"Document(source={src}, len={len(self.content)})"


class DocumentLoader:
    """Load and parse supported document formats into text chunks."""

    SUPPORTED_EXTENSIONS = {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".jsonl",
        ".pdf",
        ".docx",
        ".doc",
        ".xlsx",
        ".xls",
        ".html",
        ".htm",
        ".pptx",
        ".epub",
    }

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap 必须小于 chunk_size。")
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def load_file(self, file_path: str) -> list[Document]:
        """Load a single file and return document chunks."""
        path = Path(file_path)
        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        text = self._extract_text(path, ext)
        chunks = self._split_text(text)
        documents: list[Document] = []
        for idx, chunk in enumerate(chunks):
            content_hash = hashlib.sha256(chunk.encode("utf-8")).hexdigest()
            documents.append(
                Document(
                    content=chunk,
                    metadata={
                        "source": str(path),
                        "filename": path.name,
                        "ext": ext,
                        "chunk_index": idx,
                        "content_hash": content_hash,
                    },
                )
            )
        return documents

    def load_directory(self, dir_path: str, recursive: bool = True) -> list[Document]:
        """Load all supported files from a directory."""
        path = Path(dir_path)
        if not path.is_dir():
            raise ValueError(f"Not a directory: {dir_path}")

        docs: list[Document] = []
        pattern = "**/*" if recursive else "*"
        for file_path in path.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                try:
                    docs.extend(self.load_file(str(file_path)))
                except Exception as exc:
                    print(f"Warning: Failed to load {file_path}: {exc}")
        return docs

    def _extract_text(self, path: Path, ext: str) -> str:
        if ext in (".txt", ".md", ".json", ".jsonl"):
            return self._load_text(path)
        if ext == ".pdf":
            return self._load_pdf(path)
        if ext in (".docx", ".doc"):
            return self._load_docx(path)
        if ext in (".xlsx", ".xls"):
            return self._load_excel(path)
        if ext == ".csv":
            return self._load_csv(path)
        if ext in (".html", ".htm"):
            return self._load_html(path)
        if ext == ".pptx":
            return self._load_pptx(path)
        if ext == ".epub":
            return self._load_epub(path)
        return self._load_text(path)

    @staticmethod
    def _load_text(path: Path) -> str:
        return path.read_text(encoding="utf-8")

    @staticmethod
    def _load_pdf(path: Path) -> str:
        try:
            import pymupdf

            doc = pymupdf.open(str(path))
            try:
                return "\n\n".join(page.get_text() for page in doc)
            finally:
                doc.close()
        except ImportError:
            pass

        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        except ImportError as exc:
            raise ImportError("PDF parsing requires pypdf. Install: pip install -r requirements/rag.txt") from exc

    @staticmethod
    def _load_docx(path: Path) -> str:
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise ImportError("DOCX parsing requires python-docx. Install: pip install -r requirements/rag.txt") from exc

        doc = DocxDocument(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)

    @staticmethod
    def _load_excel(path: Path) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise ImportError("Excel parsing requires openpyxl. Install: pip install -r requirements/rag.txt") from exc

        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        try:
            lines: list[str] = []
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                lines.append(f"[Sheet: {sheet_name}]")
                for row in worksheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(cells):
                        lines.append(" | ".join(cells))
            return "\n".join(lines)
        finally:
            workbook.close()

    @staticmethod
    def _load_csv(path: Path) -> str:
        rows: list[str] = []
        with open(path, "r", encoding="utf-8", newline="") as file:
            reader = csv.reader(file)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows)

    @staticmethod
    def _load_html(path: Path) -> str:
        html = path.read_text(encoding="utf-8")
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            for element in soup(["script", "style"]):
                element.decompose()
            return soup.get_text(separator="\n", strip=True)
        except ImportError:
            text = re.sub(r"<[^>]+>", " ", html)
            return re.sub(r"\s+", " ", text).strip()

    @staticmethod
    def _load_pptx(path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError("PPTX parsing requires python-pptx. Install: pip install -r requirements/rag.txt") from exc

        presentation = Presentation(str(path))
        texts: list[str] = []
        for idx, slide in enumerate(presentation.slides, 1):
            slide_texts = [f"[Slide {idx}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            texts.append("\n".join(slide_texts))
        return "\n\n".join(texts)

    @staticmethod
    def _load_epub(path: Path) -> str:
        try:
            import ebooklib
            from bs4 import BeautifulSoup
            from ebooklib import epub
        except ImportError as exc:
            raise ImportError("EPUB parsing requires ebooklib and beautifulsoup4. Install: pip install -r requirements/rag.txt") from exc

        book = epub.read_epub(str(path))
        texts: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text:
                texts.append(text)
        return "\n\n".join(texts)

    def _split_text(self, text: str) -> list[str]:
        """Split by paragraphs and Chinese punctuation before fixed-size fallback."""
        text = re.sub(r"\r\n?", "\n", text).strip()
        if not text:
            return []

        units: list[str] = []
        for paragraph in re.split(r"\n{2,}", text):
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            units.extend(
                part.strip()
                for part in re.split(r"(?<=[。！？!?；;])\s*|\n+", paragraph)
                if part.strip()
            )

        chunks: list[str] = []
        current = ""
        for unit in units:
            candidate = f"{current}\n{unit}".strip() if current else unit
            if len(candidate) <= self.chunk_size:
                current = candidate
                continue
            if current:
                chunks.extend(self._force_split(current))
                overlap = current[-self.chunk_overlap :] if self.chunk_overlap > 0 else ""
                current = f"{overlap}{unit}"
            else:
                chunks.extend(self._force_split(unit))
                current = ""
        if current:
            chunks.extend(self._force_split(current))
        return [chunk for chunk in chunks if chunk]

    def _force_split(self, text: str) -> list[str]:
        if len(text) <= self.chunk_size:
            return [text.strip()]
        step = max(1, self.chunk_size - self.chunk_overlap)
        return [
            text[idx : idx + self.chunk_size].strip()
            for idx in range(0, len(text), step)
            if text[idx : idx + self.chunk_size].strip()
        ]


def get_supported_extensions() -> list[str]:
    """Return list of supported file extensions."""
    return sorted(DocumentLoader.SUPPORTED_EXTENSIONS)
